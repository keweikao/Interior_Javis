#!/usr/bin/env python3
"""
Q-Check Pitch Deck — 12 slides, 15 min (with live demo)
Warm interior design palette, all Chinese terminology
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import math

# ── Colors — warm interior design palette ──
BLACK = RGBColor(0x2B, 0x25, 0x22)
DARK = RGBColor(0x3D, 0x35, 0x30)
WHITE = RGBColor(0xFA, 0xF5, 0xEF)
MID_GRAY = RGBColor(0xA0, 0x96, 0x8E)
ACCENT = RGBColor(0xC4, 0x97, 0x56)
RED = RGBColor(0xC0, 0x5C, 0x4C)
ORANGE = RGBColor(0xD4, 0x8B, 0x4A)
GOLD = RGBColor(0xC4, 0x97, 0x56)
BLUE = RGBColor(0x7E, 0xA0, 0x9F)
GREEN = RGBColor(0x8B, 0xA8, 0x7E)
DARK_RED_BG = RGBColor(0x3A, 0x2A, 0x27)
DARK_GREEN_BG = RGBColor(0x2E, 0x38, 0x2E)

# Competition slide colors
AMBER = RGBColor(0xEF, 0x9F, 0x27)
AMBER_LIGHT = RGBColor(0xFD, 0xF3, 0xE0)
TEAL_C = RGBColor(0x5D, 0xCA, 0xA5)
TEAL_C_LIGHT = RGBColor(0xE8, 0xF8, 0xF0)
PURPLE_C = RGBColor(0xAF, 0xA9, 0xEC)
PURPLE_C_LIGHT = RGBColor(0xF0, 0xEE, 0xFB)
SKY_C = RGBColor(0x85, 0xB7, 0xEB)
SKY_C_LIGHT = RGBColor(0xEB, 0xF3, 0xFB)
GRAY_BORDER = RGBColor(0xC8, 0xC2, 0xBA)
GRAY_LIGHT = RGBColor(0xF5, 0xF3, 0xF0)
COMP_TEXT = RGBColor(0x2B, 0x25, 0x22)
COMP_MID = RGBColor(0x8A, 0x82, 0x7A)
AXIS_COLOR = RGBColor(0xD0, 0xC8, 0xC0)

TOTAL = 12

# ── Fill-in values ──
REVIEW_BEFORE = "30"
REVIEW_AFTER = "10"
TEST_REPORTS = "5"
TEST_WEEKS = "2"
TEST_CASES = "3"
MONTHLY_PRICE = "990"
NETWORK_COUNT = "50+"
INTERVIEW_COUNT = "5"
BETA_COUNT = "3-5"
COMPETITOR_NAME = "___（待填）"
MVP_WEEKS = "4-5"
EMAIL = "hello@qcheck.tw"
DEMO_URL = "https://q-check-demo.pages.dev"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def bg(slide, c=BLACK):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def tx(slide, l, t, w, h, text, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = "Microsoft JhengHei"; p.alignment = a
    return tf

def ap(tf, text, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT, sp=Pt(6)):
    p = tf.add_paragraph(); p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = "Microsoft JhengHei"; p.alignment = a
    if sp: p.space_before = sp
    return p

def bar(slide, l, t, w=Inches(0.8), h=Inches(0.06)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

def card(slide, l, t, w, h, c=DARK):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def divider(slide, t):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), t, Inches(11.7), Pt(1))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0x40, 0x40, 0x55); s.line.fill.background()

def sn(slide, n):
    tx(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
       f"{n}/{TOTAL}", sz=11, c=MID_GRAY, a=PP_ALIGN.RIGHT)

def lbl(slide, text):
    tx(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.5), text, sz=14, c=ACCENT, b=True)
    bar(slide, Inches(0.8), Inches(1.1), Inches(0.6))

def heading(slide, text):
    tx(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.8), text, sz=28, b=True)


# ════════════════════════════════════════
# 1 — COVER
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, Inches(0.8), Inches(2.2), Inches(1.2), Inches(0.08))
tx(s, Inches(0.8), Inches(2.5), Inches(10), Inches(1.2), "Q-Check", sz=60, b=True, c=ACCENT)

# Title: 室內裝修公司的 + 風險管理決策系統（粗體）
tb = s.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(10), Inches(0.8))
ttf = tb.text_frame; ttf.word_wrap = True
p = ttf.paragraphs[0]
run1 = p.add_run(); run1.text = "室內裝修公司的"
run1.font.size = Pt(32); run1.font.color.rgb = WHITE; run1.font.name = "Microsoft JhengHei"
run2 = p.add_run(); run2.text = "風險管理決策系統"
run2.font.size = Pt(32); run2.font.bold = True; run2.font.color.rgb = WHITE; run2.font.name = "Microsoft JhengHei"

tx(s, Inches(0.8), Inches(4.5), Inches(10), Inches(0.6),
   "讓報價階段的錯誤，不會留到施工階段。", sz=20, c=MID_GRAY)
tx(s, Inches(0.8), Inches(5.8), Inches(10), Inches(0.4),
   "第六組：呂賀云  /  高克瑋", sz=14, c=MID_GRAY)
tx(s, Inches(0.8), Inches(6.2), Inches(10), Inches(0.4),
   "業師：張佳欽    2026-03-24", sz=14, c=MID_GRAY)


# ════════════════════════════════════════
# 2 — PAIN + WHY NO ONE SOLVED IT
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); sn(s, 2)
lbl(s, "問題")

tx(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
   "入行 2 年的設計師接了 30 坪中古屋翻新案", sz=26, b=True)

tf = tx(s, Inches(0.8), Inches(2.3), Inches(7), Inches(2), "", sz=20)
ap(tf, "5 樓無電梯、屋齡 25 年、業主預算 180 萬", sz=20, c=WHITE, sp=Pt(12))
ap(tf, "報價漏了防水重做、垃圾清運、配電箱更換", sz=20, c=ORANGE, sp=Pt(12))
ap(tf, "施工後追加 35 萬 — 業主崩潰、公司賠錢", sz=20, c=RED, b=True, sp=Pt(12))

pains = [
    ("認知落差", "報價寫了但雙方對「含什麼」理解不同，容易產生糾紛"),
    ("漏項", "清運、防水、五金等項目常遺漏，直接侵蝕毛利"),
    ("現場條件低估", "無電梯搬運、老屋牆面狀況，追加費用更傷"),
    ("預算不合理", "接了預算做不了的案子，整案利潤全失"),
]
for i, (t, d) in enumerate(pains):
    y = Inches(2.3) + Inches(i * 0.7)
    tx(s, Inches(8.5), y, Inches(4), Inches(0.35), t, sz=13, c=ACCENT, b=True)
    tx(s, Inches(8.5), y + Inches(0.32), Inches(4), Inches(0.35), d, sz=10, c=MID_GRAY)

tx(s, Inches(0.8), Inches(4.6), Inches(7), Inches(0.5),
   "報價階段修復 vs 施工階段修復，代價差 10 到 20 倍", sz=16, c=GOLD, b=True)

divider(s, Inches(5.2))
tx(s, Inches(0.8), Inches(5.4), Inches(11), Inches(0.4),
   "為什麼現有方案都解決不了？", sz=16, c=ACCENT, b=True)

fails = [
    ("Excel 模板", "有格式但沒有風險檢查，知識留在人腦裡"),
    ("口頭傳承", "資深沒時間一對一教，帶 3 個人就是極限"),
    ("工程管理軟體", "管施工不管報價，對設計公司太重太貴"),
]
for i, (n, r) in enumerate(fails):
    x = Inches(0.8) + Inches(i * 4.1)
    card(s, x, Inches(5.9), Inches(2.9), Inches(1.1), DARK)
    tx(s, x + Inches(0.15), Inches(5.95), Inches(2.6), Inches(0.35),
       n, sz=13, c=RED, b=True)
    tx(s, x + Inches(0.15), Inches(6.35), Inches(2.6), Inches(0.5),
       r, sz=11, c=MID_GRAY)


# ════════════════════════════════════════
# 3 — PRODUCT: 4 Layers + 專業判斷
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); sn(s, 3)
lbl(s, "Q-Check 怎麼運作")

tx(s, Inches(0.8), Inches(1.4), Inches(11), Inches(0.6),
   "培養設計師，不是取代設計師。", sz=26, b=True)
tx(s, Inches(0.8), Inches(2.0), Inches(11), Inches(0.4),
   "— 呂賀云（共同創辦人 / 10 年資深設計師）", sz=14, c=MID_GRAY)

steps = [
    ("01", "預算快篩", "這案子值不值得接？", "輸入坪數、屋齡、預算，系統判斷預算是否合理", ACCENT),
    ("02", "結構化報價", "該報哪些項目？", "選案件類型自動帶入工項骨架，每項填寫含什麼、不含什麼", BLUE),
    ("03", "即時風險提醒", "有沒有漏掉什麼？", "根據 25 條以上工序連動規則，在旁邊即時提示風險和漏項", ORANGE),
    ("04", "完成確認", "可以安心送出了嗎？", "逐項確認後產出 PDF 報價單和內部風險報告", GREEN),
]
for i, (num, t, q, d, c) in enumerate(steps):
    y = Inches(2.6) + Inches(i * 0.95)
    card(s, Inches(0.8), y, Inches(7.4), Inches(0.75), DARK)
    tx(s, Inches(1.0), y + Inches(0.1), Inches(0.6), Inches(0.5),
       num, sz=18, c=c, b=True)
    tx(s, Inches(1.6), y + Inches(0.1), Inches(1.8), Inches(0.3),
       t, sz=14, b=True)
    tx(s, Inches(3.4), y + Inches(0.1), Inches(2.5), Inches(0.3),
       q, sz=12, c=c)
    tx(s, Inches(1.6), y + Inches(0.42), Inches(6.3), Inches(0.3),
       d, sz=11, c=MID_GRAY)

card(s, Inches(8.5), Inches(2.6), Inches(4.0), Inches(4.2), DARK)
tx(s, Inches(8.8), Inches(2.75), Inches(3.5), Inches(0.4),
   "「專業判斷」機制", sz=15, c=GOLD, b=True)

tf = tx(s, Inches(8.8), Inches(3.2), Inches(3.5), Inches(3.2), "", sz=13)
ap(tf, "系統提醒：「這個單價低於市場行情」", sz=13, c=ORANGE, sp=Pt(8))
ap(tf, "", sz=8, c=MID_GRAY, sp=Pt(2))
ap(tf, "設計師可以選擇：", sz=13, c=WHITE, b=True, sp=Pt(6))
ap(tf, "  接受建議，調整單價", sz=12, c=WHITE, sp=Pt(3))
ap(tf, "  或者行使專業判斷，寫下原因", sz=12, c=WHITE, sp=Pt(3))
ap(tf, "  例如：「合作多年工班，有特殊價格」", sz=11, c=MID_GRAY, sp=Pt(2))
ap(tf, "", sz=8, c=MID_GRAY, sp=Pt(4))
ap(tf, "系統記錄這個決定，不再重複提醒", sz=12, c=ACCENT, sp=Pt(4))
ap(tf, "資深覆核時可以看到所有判斷紀錄", sz=12, c=ACCENT, sp=Pt(3))

tx(s, Inches(8.8), Inches(6.0), Inches(3.5), Inches(0.7),
   "這讓系統分得出「故意的」和「不小心的」", sz=13, c=GOLD)


# ════════════════════════════════════════
# 4 — LIVE DEMO (minimal)
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tx(s, Inches(0.8), Inches(2.6), Inches(11.5), Inches(1.2),
   "Live Demo", sz=52, b=True, c=ACCENT, a=PP_ALIGN.CENTER)
tx(s, Inches(0.8), Inches(3.9), Inches(11.5), Inches(0.6),
   "用剛才的案子走一遍：30 坪中古屋、屋齡 25 年、預算 180 萬", sz=22, c=MID_GRAY, a=PP_ALIGN.CENTER)
tx(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5),
   DEMO_URL, sz=18, c=ACCENT, a=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# 5 — VALUE PROPOSITION CANVAS (3 roles, classic VPC per role)
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); sn(s, 5)
lbl(s, "價值主張畫布")
heading(s, "三個角色的價值互相增強")

vpc = [
    ("年輕設計師", BLUE,
     # 顧客任務
     ["獨立完成一份完整報價單", "不漏項、不寫模糊", "讓報價通過資深覆核"],
     # 痛點
     ["不知道自己漏了什麼", "要等資深空下來才敢發報價", "報價常被退回修改好幾次"],
     # 期望收穫
     ["報價品質提升", "不用等人就敢發", "客戶覺得很專業"],
     # 產品與服務
     ["工項骨架一鍵帶入", "含/不含欄位引導", "即時風險提醒附解釋"],
     # 怎麼減少痛苦
     ["骨架帶入不會遺漏", "提醒附「為什麼」是在學習", "預算快篩不接錯案"],
     # 怎麼創造收穫
     ["報價一次通過覆核", "累積知識漸漸不需提醒", "專業判斷紀錄展現成長"]),
    ("資深設計師", GOLD,
     ["覆核年輕人的報價", "同時管多個進行中案件", "把經驗傳下去"],
     ["大量時間花在覆核", "同樣的錯教了很多次", "自己的案在延遲"],
     ["覆核時間縮短", "經驗變成公司資產", "看得出年輕人成長"],
     ["風險檢核報告", "專業判斷清單", "系統已檢查項目清單"],
     ["風險報告過濾規則性問題", "只看需要經驗判斷的部分", "判斷紀錄分辨故意或疏忽"],
     ["省下時間經營客戶接新案", "經驗透過規則庫傳承", "看判斷紀錄掌握成長速度"]),
    ("老闆", GREEN,
     ["用現有人力接更多案", "控制追加風險和毛利", "培養年輕設計師獨立"],
     ["人力有限案量受限", "年輕人要很久才獨立", "追加費用侵蝕毛利"],
     ["同樣人力接更多案", "年輕設計師更快獨立", "品質不依賴特定個人"],
     ["預算可行性快篩", "品質標準化系統", "PDF 報價單匯出"],
     ["預算快篩避免接錯案", "品質標準化不靠個人", "報價格式統一"],
     ["資深時間釋放接更多案", "年輕人更快能獨立作業", "知識留在公司不隨人走"]),
]

for i, (role, color, jobs, pains_list, gains, products, relievers, creators) in enumerate(vpc):
    x = Inches(0.5) + Inches(i * 4.2)
    cw = Inches(4.0)

    # Role header
    card(s, x, Inches(2.25), cw, Inches(0.4), color)
    tx(s, x + Inches(0.15), Inches(2.28), cw - Inches(0.3), Inches(0.3),
       role, sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)

    # Left half: Value Map (產品/減痛/創穫)
    half_w = Inches(1.9)
    card(s, x, Inches(2.75), half_w, Inches(4.2), DARK)

    ty = Inches(2.8)
    tx(s, x + Inches(0.08), ty, half_w - Inches(0.15), Inches(0.2),
       "產品與服務", sz=8, b=True, c=ACCENT)
    tf = tx(s, x + Inches(0.08), ty + Inches(0.2), half_w - Inches(0.15), Inches(0.9), "", sz=9)
    for p in products:
        ap(tf, p, sz=9, c=WHITE, sp=Pt(1))

    ty2 = Inches(3.9)
    tx(s, x + Inches(0.08), ty2, half_w - Inches(0.15), Inches(0.2),
       "怎麼減少痛苦", sz=8, b=True, c=RED)
    tf = tx(s, x + Inches(0.08), ty2 + Inches(0.2), half_w - Inches(0.15), Inches(0.9), "", sz=9)
    for p in relievers:
        ap(tf, p, sz=9, c=WHITE, sp=Pt(1))

    ty3 = Inches(5.0)
    tx(s, x + Inches(0.08), ty3, half_w - Inches(0.15), Inches(0.2),
       "怎麼創造收穫", sz=8, b=True, c=GREEN)
    tf = tx(s, x + Inches(0.08), ty3 + Inches(0.2), half_w - Inches(0.15), Inches(0.9), "", sz=9)
    for p in creators:
        ap(tf, p, sz=9, c=color, sp=Pt(1))

    # Right half: Customer Profile (任務/痛點/期望)
    rx = x + half_w + Inches(0.05)
    card(s, rx, Inches(2.75), half_w, Inches(4.2), DARK)

    tx(s, rx + Inches(0.08), ty, half_w - Inches(0.15), Inches(0.2),
       "顧客任務", sz=8, b=True, c=ACCENT)
    tf = tx(s, rx + Inches(0.08), ty + Inches(0.2), half_w - Inches(0.15), Inches(0.9), "", sz=9)
    for j in jobs:
        ap(tf, j, sz=9, c=WHITE, sp=Pt(1))

    tx(s, rx + Inches(0.08), ty2, half_w - Inches(0.15), Inches(0.2),
       "痛點", sz=8, b=True, c=RED)
    tf = tx(s, rx + Inches(0.08), ty2 + Inches(0.2), half_w - Inches(0.15), Inches(0.9), "", sz=9)
    for p in pains_list:
        ap(tf, p, sz=9, c=WHITE, sp=Pt(1))

    tx(s, rx + Inches(0.08), ty3, half_w - Inches(0.15), Inches(0.2),
       "期望收穫", sz=8, b=True, c=GREEN)
    tf = tx(s, rx + Inches(0.08), ty3 + Inches(0.2), half_w - Inches(0.15), Inches(0.9), "", sz=9)
    for g in gains:
        ap(tf, g, sz=9, c=color, sp=Pt(1))

# Flywheel
tx(s, Inches(0.5), Inches(7.0), Inches(12.5), Inches(0.35),
   "三者互相增強：年輕人底氣提升 → 資深覆核減少 → 老闆產能提升", sz=12, c=ACCENT, b=True, a=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# 6 — MARKET + WHY NOW (no competition)
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); sn(s, 6)
lbl(s, "市場機會")
heading(s, "台灣裝修年產值 5,500 億，風險管理是空白地帶")

# Funnel (left)
funnel = [
    ("總潛在市場", "全台設計公司約 17,000 家", Inches(7)),
    ("可服務市場", "3 人以上中型公司約 3,000-4,000 家", Inches(5.5)),
    ("首年目標", "人脈能觸及 50-100 家", Inches(4)),
]
for i, (lb_t, desc, w) in enumerate(funnel):
    y = Inches(2.3) + Inches(i * 0.9)
    x = Inches(0.8) + (Inches(7) - w) / 2
    colors = [RGBColor(0x4A, 0x3F, 0x38), RGBColor(0x5A, 0x4E, 0x45), ACCENT]
    card(s, x, y, w, Inches(0.7), colors[i])
    tx(s, x + Inches(0.2), y + Inches(0.1), Inches(1.8), Inches(0.35),
       lb_t, sz=13, b=True, c=WHITE)
    tx(s, x + Inches(2.2), y + Inches(0.12), w - Inches(2.5), Inches(0.35),
       desc, sz=13, c=WHITE)

# Why Now (right side)
card(s, Inches(8.3), Inches(2.3), Inches(4.2), Inches(4.5), DARK)
tx(s, Inches(8.6), Inches(2.4), Inches(3.8), Inches(0.35),
   "為什麼是現在", sz=14, b=True, c=ACCENT)

reasons = [
    ("老屋翻新潮", "500 萬戶老屋加上政府 50 億補助\n案量增加但產能不變"),
    ("報價失準代價更高", "材料成本上漲 6.8%\n報價少算就是直接賠錢"),
    ("經驗傳承斷層", "資深設計師退休潮\n年輕人學不到也沒人教"),
    ("沒人做這件事", "目前沒人從風險管理切入\n要趁空白搶先建立用戶"),
]
for i, (t, d) in enumerate(reasons):
    y = Inches(2.9) + Inches(i * 0.95)
    tx(s, Inches(8.6), y, Inches(3.8), Inches(0.3), t, sz=12, c=ACCENT, b=True)
    tx(s, Inches(8.6), y + Inches(0.3), Inches(3.8), Inches(0.55), d, sz=11, c=WHITE)


# ════════════════════════════════════════
# 7 — POCD
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); sn(s, 7)
lbl(s, "POCD 分析")

qw = Inches(5.85); qh = Inches(2.7); g = Inches(0.1)

# People (no advisor)
card(s, Inches(0.8), Inches(1.4), qw, qh, DARK)
tx(s, Inches(1.1), Inches(1.5), Inches(3), Inches(0.35), "People 團隊", sz=16, b=True, c=ACCENT)
tf = tx(s, Inches(1.1), Inches(1.9), qw - Inches(0.5), Inches(2), "", sz=13)
ap(tf, "呂賀云：10 年以上設計師，她本人就是目標用戶", sz=13, c=WHITE, sp=Pt(3))
ap(tf, "高克瑋：10 年銷售營運經驗，iCHEF 銷售流程 AI 導入經驗", sz=13, c=WHITE, sp=Pt(3))
ap(tf, "擅長把專家知識變成可執行的系統", sz=13, c=WHITE, sp=Pt(3))
ap(tf, "共同創辦人就是目標用戶，對市場的理解極深", sz=13, c=ACCENT, b=True, sp=Pt(6))

# Opportunity
ox = Inches(0.8) + qw + g
card(s, ox, Inches(1.4), qw, qh, DARK)
tx(s, ox + Inches(0.2), Inches(1.5), Inches(3), Inches(0.35), "Opportunity 機會", sz=16, b=True, c=ACCENT)
tf = tx(s, ox + Inches(0.2), Inches(1.9), qw - Inches(0.5), Inches(2), "", sz=13)
ap(tf, "全台約 17,000 家設計公司", sz=13, c=WHITE, sp=Pt(3))
ap(tf, "目標客群：3 人以上中型公司約 3,000-4,000 家", sz=13, c=WHITE, sp=Pt(3))
ap(tf, f"賀云的 13 家分公司加上 {NETWORK_COUNT} 家直接人脈", sz=13, c=WHITE, sp=Pt(3))
ap(tf, "從風險管理角度切入，目前沒有直接競品", sz=13, c=GOLD, b=True, sp=Pt(6))

# Context (no window period, PLG adoption)
cy = Inches(1.4) + qh + g
card(s, Inches(0.8), cy, qw, qh, DARK)
tx(s, Inches(1.1), cy + Inches(0.08), Inches(3), Inches(0.35), "Context 環境", sz=16, b=True, c=ACCENT)
tf = tx(s, Inches(1.1), cy + Inches(0.45), qw - Inches(0.5), Inches(2), "", sz=12)
ap(tf, "有利：老屋翻新潮、裝潢糾紛頻繁、部署成本趨近零", sz=12, c=GREEN, sp=Pt(3))
ap(tf, "不利：設計業數位化程度低、各公司格式不統一", sz=12, c=RED, sp=Pt(3))
ap(tf, "擴散方式：設計師自己覺得好用 → 推薦同事 → 公司導入", sz=12, c=WHITE, sp=Pt(3))
ap(tf, "類似 Figma 的擴散路徑，從使用者端長出來", sz=12, c=WHITE, sp=Pt(3))

# Deal (ARPU, 3-tier pricing, no MRR)
card(s, ox, cy, qw, qh, DARK)
tx(s, ox + Inches(0.2), cy + Inches(0.08), Inches(3), Inches(0.35), "Deal 商業模式", sz=16, b=True, c=ACCENT)
tf = tx(s, ox + Inches(0.2), cy + Inches(0.45), qw - Inches(0.5), Inches(2), "", sz=12)
ap(tf, "免費版：個人設計師，註冊就能用", sz=12, c=WHITE, sp=Pt(3))
ap(tf, f"進階版：NT$ {MONTHLY_PRICE}/月，分享給同事、更多範本", sz=12, c=WHITE, sp=Pt(3))
ap(tf, "企業版：另議，覆核報告、多人管理、品牌 PDF", sz=12, c=WHITE, sp=Pt(3))
ap(tf, f"擴張路徑：13 家分公司 → {NETWORK_COUNT} 家人脈 → 口碑擴散", sz=12, c=WHITE, sp=Pt(6))
ap(tf, "獲客成本趨近零   部署成本趨近零", sz=12, c=ACCENT, b=True, sp=Pt(4))


# ════════════════════════════════════════
# 8 — COMPETITION MAP (warm dark bg, matching deck style)
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); sn(s, 8)
lbl(s, "競爭格局")
heading(s, "每個玩家賣的是什麼承諾？")

# Warm dark card fills for competition
COMP_CARD = RGBColor(0x4A, 0x42, 0x3C)
COMP_QCHECK = RGBColor(0x4D, 0x3D, 0x28)  # warm amber-tinted dark

CHART_L = 1.8; CHART_R = 11.8; CHART_T = 2.0; CHART_B = 6.6
CHART_CX = (CHART_L + CHART_R) / 2; CHART_CY = (CHART_T + CHART_B) / 2

def cxp(pct): return CHART_L + (CHART_R - CHART_L) * pct / 100
def cyp(pct): return CHART_B - (CHART_B - CHART_T) * pct / 100

# Axes (warm subtle)
WARM_AXIS = RGBColor(0x5A, 0x52, 0x4A)
vl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(CHART_CX) - Pt(0.5), Inches(CHART_T), Pt(1), Inches(CHART_B - CHART_T))
vl.fill.solid(); vl.fill.fore_color.rgb = WARM_AXIS; vl.line.fill.background()

hl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(CHART_L), Inches(CHART_CY) - Pt(0.5), Inches(CHART_R - CHART_L), Pt(1))
hl.fill.solid(); hl.fill.fore_color.rgb = WARM_AXIS; hl.line.fill.background()

def ax_lbl(x, y, text, align=PP_ALIGN.CENTER):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(2), Inches(0.3))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(11); p.font.color.rgb = MID_GRAY; p.font.name = "Microsoft JhengHei"
    p.alignment = align

ax_lbl(CHART_L - 0.2, CHART_CY + 0.08, "通用工具", PP_ALIGN.LEFT)
ax_lbl(CHART_R - 1.8, CHART_CY + 0.08, "室內設計專用", PP_ALIGN.RIGHT)
ax_lbl(CHART_CX - 1.0, CHART_T - 0.32, "事前預防")
ax_lbl(CHART_CX - 1.0, CHART_B + 0.05, "事後管理")

def comp_card(x_pct, y_pct, name, promise, detail,
              accent_c=MID_GRAY, fill_c=COMP_CARD,
              cw=2.3, ch=1.1, bw=Pt(1), nsz=13, psz=10, dsz=9):
    lx = cxp(x_pct) - cw / 2; ly = cyp(y_pct) - ch / 2
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(lx), Inches(ly), Inches(cw), Inches(ch))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill_c
    sh.line.color.rgb = accent_c; sh.line.width = bw

    tx(s, Inches(lx+0.12), Inches(ly+0.08), Inches(cw-0.24), Inches(0.28),
       name, sz=nsz, b=True, c=WHITE)
    tx(s, Inches(lx+0.12), Inches(ly+0.38), Inches(cw-0.24), Inches(0.28),
       f"「{promise}」", sz=psz, c=accent_c)
    if detail:
        tx(s, Inches(lx+0.12), Inches(ly+0.68), Inches(cw-0.24), Inches(0.35),
           detail, sz=dsz, c=MID_GRAY)
    return lx, ly, cw, ch

# Q-Check (largest, amber accent)
qlx, qly, qlw, qlh = comp_card(78, 82, "Q-Check", "這份報價有沒有埋地雷？",
    "風險檢核 · 含/不含 · 專業判斷",
    accent_c=ACCENT, fill_c=COMP_QCHECK, cw=2.7, ch=1.2, bw=Pt(2.5), nsz=15, psz=11, dsz=10)

# Dashed highlight
dm = 0.15
ds = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(qlx-dm), Inches(qly-dm), Inches(qlw+dm*2), Inches(qlh+dm*2+0.32))
ds.fill.background(); ds.line.color.rgb = ACCENT; ds.line.width = Pt(1.5)
prd = etree.SubElement(ds.line._ln, qn('a:prstDash')); prd.set('val', 'dash')

tx(s, Inches(qlx-dm), Inches(qly+qlh+0.03), Inches(qlw+dm*2), Inches(0.28),
   "唯一的事前預防 × 設計專用", sz=10, b=True, c=ACCENT, a=PP_ALIGN.CENTER)

# Other cards
comp_card(73, 40, "易裝修 EZid", "案子跑更快、團隊不亂",
    "$1,200/月起 · 報價生成 · 流程管理", accent_c=BLUE)
comp_card(25, 52, "Excel 報價模板", "我最懂自己的案子",
    "免費 · 70-80% 設計師在用 · 零檢核零傳承")
comp_card(25, 30, "LINE 群組", "傳個訊息就搞定了",
    "免費 · 90% 在用 · 即時但無法追溯")
comp_card(75, 17, "100室內設計", "幫你找到下一個客戶",
    "消費者媒合平台 · 2,000+ 設計公司", accent_c=RGBColor(0x9A, 0x94, 0xD0))
comp_card(25, 75, "元欣估價系統", "估價標準化不出錯",
    "單機版 · 通用工程 · 不懂裝修工法", accent_c=BLUE)
comp_card(20, 13, "Notion / Trello", "什麼都能管 = 什麼都不深", "",
    cw=2.0, ch=0.85, nsz=11, psz=9)


# ════════════════════════════════════════
# 9 — LEAN CANVAS
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); sn(s, 9)
tx(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.35),
   "精實畫布", sz=14, c=ACCENT, b=True)
bar(s, Inches(0.8), Inches(0.68), Inches(0.6))

r1y = Inches(0.9); r1h = Inches(3.2); cw2 = Inches(2.35); g2 = Inches(0.06)

# Problem
x = Inches(0.8)
card(s, x, r1y, cw2, r1h, DARK)
tx(s, x+Inches(0.1), r1y+Inches(0.05), cw2-Inches(0.2), Inches(0.22), "問題", sz=11, b=True, c=RED)
tf = tx(s, x+Inches(0.1), r1y+Inches(0.28), cw2-Inches(0.2), Inches(1.3), "", sz=10)
ap(tf, "1. 報價寫了但雙方對「含什麼」理解不同", sz=10, c=WHITE, sp=Pt(2))
ap(tf, "2. 保護、清運、防水等項目常被遺漏", sz=10, c=WHITE, sp=Pt(2))
ap(tf, "3. 資深設計師大量時間花在覆核", sz=10, c=WHITE, sp=Pt(2))
tx(s, x+Inches(0.1), r1y+Inches(1.6), cw2-Inches(0.2), Inches(0.2), "現有替代方案", sz=9, b=True, c=MID_GRAY)
tf2 = tx(s, x+Inches(0.1), r1y+Inches(1.8), cw2-Inches(0.2), Inches(1), "", sz=9)
ap(tf2, "Excel 模板、口頭傳承、出事後檢討", sz=9, c=MID_GRAY, sp=Pt(1))

# Solution
x += cw2 + g2
card(s, x, r1y, cw2, r1h, DARK)
tx(s, x+Inches(0.1), r1y+Inches(0.05), cw2-Inches(0.2), Inches(0.22), "解決方案", sz=11, b=True, c=ACCENT)
tf = tx(s, x+Inches(0.1), r1y+Inches(0.28), cw2-Inches(0.2), Inches(2.5), "", sz=10)
ap(tf, "1. 含/不含欄位強制寫清楚", sz=10, c=WHITE, sp=Pt(2))
ap(tf, "2. 25 條以上工序連動規則即時提醒", sz=10, c=WHITE, sp=Pt(2))
ap(tf, "3. 預算可行性快篩", sz=10, c=WHITE, sp=Pt(2))
ap(tf, "4. 風險報告讓資深只看需判斷的部分", sz=10, c=WHITE, sp=Pt(2))
ap(tf, "5. 專業判斷機制（覆寫加註原因）", sz=10, c=WHITE, sp=Pt(2))

# UVP
x += cw2 + g2
card(s, x, r1y, cw2, r1h, DARK)
tx(s, x+Inches(0.1), r1y+Inches(0.05), cw2-Inches(0.2), Inches(0.22), "獨特價值主張", sz=11, b=True, c=GOLD)
tf = tx(s, x+Inches(0.1), r1y+Inches(0.28), cw2-Inches(0.2), Inches(1.2), "", sz=10)
ap(tf, "讓報價階段的錯誤，不留到施工階段", sz=10, c=WHITE, b=True, sp=Pt(2))
ap(tf, "不是報價工具，是風險管理系統", sz=9, c=MID_GRAY, sp=Pt(4))
ap(tf, "培養設計師的底氣，不是取代設計師", sz=9, c=MID_GRAY, sp=Pt(2))
tx(s, x+Inches(0.1), r1y+Inches(1.4), cw2-Inches(0.2), Inches(0.2), "不公平競爭優勢", sz=9, b=True, c=GOLD)
tf2 = tx(s, x+Inches(0.1), r1y+Inches(1.6), cw2-Inches(0.2), Inches(1.3), "", sz=10)
ap(tf2, "賀云本人就是目標用戶", sz=10, c=WHITE, sp=Pt(1))
ap(tf2, "10 年實戰知識已編碼為規則", sz=10, c=WHITE, sp=Pt(1))
ap(tf2, "用戶使用就會產生乾淨數據", sz=10, c=WHITE, sp=Pt(1))
ap(tf2, "早期靠產業知識壁壘", sz=10, c=WHITE, sp=Pt(1))
ap(tf2, "中期靠知識飛輪", sz=10, c=ACCENT, sp=Pt(1))

# Key Metrics
x += cw2 + g2
nw = Inches(2.0)
card(s, x, r1y, nw, r1h, DARK)
tx(s, x+Inches(0.1), r1y+Inches(0.05), nw-Inches(0.2), Inches(0.22), "關鍵指標", sz=11, b=True, c=BLUE)
tf = tx(s, x+Inches(0.1), r1y+Inches(0.28), nw-Inches(0.2), Inches(2.5), "", sz=10)
ap(tf, "追加金額變化", sz=10, c=WHITE, sp=Pt(2))
ap(tf, "覆核時間變化", sz=10, c=WHITE, sp=Pt(2))
ap(tf, "報價來回修改次數", sz=10, c=WHITE, sp=Pt(2))
ap(tf, "每月活躍專案數", sz=10, c=WHITE, sp=Pt(2))
ap(tf, "專業判斷使用頻率", sz=10, c=WHITE, sp=Pt(2))

# Customer
x += nw + g2
card(s, x, r1y, cw2, r1h, DARK)
tx(s, x+Inches(0.1), r1y+Inches(0.05), cw2-Inches(0.2), Inches(0.22), "目標客群", sz=11, b=True, c=GREEN)
tf = tx(s, x+Inches(0.1), r1y+Inches(0.28), cw2-Inches(0.2), Inches(1.5), "", sz=10)
ap(tf, "使用者：設計師（先免費用起來）", sz=10, c=ACCENT, b=True, sp=Pt(2))
ap(tf, "付費者：進階功能的設計師或公司", sz=10, c=ACCENT, b=True, sp=Pt(2))
ap(tf, "", sz=4, c=WHITE, sp=Pt(2))
ap(tf, "3 人以上的中型設計公司", sz=10, c=WHITE, sp=Pt(1))
ap(tf, "年輕設計師需要資深覆核的團隊", sz=10, c=WHITE, sp=Pt(1))
tx(s, x+Inches(0.1), r1y+Inches(1.6), cw2-Inches(0.2), Inches(0.2), "通路", sz=9, b=True, c=MID_GRAY)
tf2 = tx(s, x+Inches(0.1), r1y+Inches(1.8), cw2-Inches(0.2), Inches(1), "", sz=9)
ap(tf2, "設計師自己覺得好用，推薦同事", sz=9, c=MID_GRAY, sp=Pt(1))
ap(tf2, "13 家分公司作為第一批種子", sz=9, c=MID_GRAY, sp=Pt(1))
ap(tf2, f"{NETWORK_COUNT} 家人脈口碑擴散", sz=9, c=MID_GRAY, sp=Pt(1))

# Row 2
r2y = r1y + r1h + Inches(0.06)
r2h = Inches(1.8)
hw = Inches(5.93)

card(s, Inches(0.8), r2y, hw, r2h, DARK)
tx(s, Inches(0.9), r2y+Inches(0.05), Inches(5), Inches(0.22), "成本結構", sz=11, b=True, c=RED)
tf = tx(s, Inches(0.9), r2y+Inches(0.28), Inches(5.5), Inches(1.3), "", sz=10)
ap(tf, "開發人力（高克瑋，機會成本）", sz=10, c=WHITE, sp=Pt(1))
ap(tf, "伺服器託管費用趨近零", sz=10, c=WHITE, sp=Pt(1))
ap(tf, "燒錢率極低，目前不需外部資金", sz=10, c=MID_GRAY, sp=Pt(1))

rx = Inches(0.8) + hw + g2
card(s, rx, r2y, hw, r2h, DARK)
tx(s, rx+Inches(0.1), r2y+Inches(0.05), Inches(5), Inches(0.22), "收入來源", sz=11, b=True, c=GREEN)
tf = tx(s, rx+Inches(0.1), r2y+Inches(0.28), Inches(5.5), Inches(1.3), "", sz=10)
ap(tf, "免費版：個人設計師，註冊就能用", sz=10, c=WHITE, sp=Pt(1))
ap(tf, f"進階版：NT$ {MONTHLY_PRICE}/月，分享同事、更多範本", sz=10, c=WHITE, sp=Pt(1))
ap(tf, "企業版：另議，覆核報告、多人管理、品牌 PDF", sz=10, c=WHITE, sp=Pt(1))
ap(tf, f"一次漏項損失 10-30 萬，月費 NT$ {MONTHLY_PRICE} 幾乎零風險", sz=10, c=GOLD, b=True, sp=Pt(3))

tx(s, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
   f"最高風險：設計師願意離開 Excel 嗎？做法：先讓設計師免費用，好用自然擴散",
   sz=12, c=ORANGE, b=True)


# ════════════════════════════════════════
# 10 — WHY US + MOAT
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); sn(s, 10)
lbl(s, "為什麼是我們")
heading(s, "共同創辦人就是目標用戶本人")

# Vivian
card(s, Inches(0.8), Inches(2.5), Inches(5.5), Inches(2.8), DARK)
tx(s, Inches(1.2), Inches(2.65), Inches(4.8), Inches(0.4),
   "呂賀云", sz=22, b=True, c=ACCENT)
tx(s, Inches(1.2), Inches(3.1), Inches(4.8), Inches(0.3),
   "共同創辦人 / 產品定義 / 設計師", sz=13, c=MID_GRAY)
tf = tx(s, Inches(1.2), Inches(3.45), Inches(4.8), Inches(1.5), "", sz=15)
ap(tf, "10 年以上室內設計實戰經驗", sz=15, c=WHITE, sp=Pt(4))
ap(tf, "她親身經歷了所有痛點，她就是用戶", sz=15, c=WHITE, sp=Pt(4))
ap(tf, "定義了產品哲學和 25 條以上的風險規則", sz=15, c=WHITE, sp=Pt(4))
ap(tf, "13 家分公司是天然的第一批用戶", sz=15, c=WHITE, sp=Pt(4))

# Stephen
card(s, Inches(6.8), Inches(2.5), Inches(5.7), Inches(2.8), DARK)
tx(s, Inches(7.2), Inches(2.65), Inches(5), Inches(0.4),
   "高克瑋", sz=22, b=True, c=ACCENT)
tx(s, Inches(7.2), Inches(3.1), Inches(5), Inches(0.3),
   "技術創辦人 / 銷售營運", sz=13, c=MID_GRAY)
tf = tx(s, Inches(7.2), Inches(3.45), Inches(5), Inches(1.5), "", sz=15)
ap(tf, "iCHEF Sales Ops Director 經驗", sz=15, c=WHITE, sp=Pt(4))
ap(tf, "擅長把專家知識變成可執行的系統", sz=15, c=WHITE, sp=Pt(4))
ap(tf, "獨立完成最小可行產品開發，不需外部資金", sz=15, c=WHITE, sp=Pt(4))

# Moat
card(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.5), DARK)
tx(s, Inches(1.2), Inches(5.7), Inches(11), Inches(0.35),
   "壁壘", sz=15, b=True, c=ACCENT)
tf = tx(s, Inches(1.2), Inches(6.05), Inches(11), Inches(0.9), "", sz=14)
ap(tf, "早期：專家知識 — 10 年經驗不是三個月能複製的", sz=14, c=WHITE, sp=Pt(3))
ap(tf, "     重點不是複製規則，是複製「知道哪些規則重要」的判斷力", sz=13, c=MID_GRAY, sp=Pt(2))
ap(tf, "中期：每一家新客戶在報價上的回饋都讓系統更好，對下一家客戶更有價值", sz=14, c=ACCENT, sp=Pt(5))


# ════════════════════════════════════════
# 11 — VALIDATION + RISK
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); sn(s, 11)
lbl(s, "驗證計畫和風險")
heading(s, "我們知道什麼還沒被驗證")

risks = [
    ("致命", "設計師願意把報價從 Excel 搬到新系統嗎？",
     f"賀云公司實測 {TEST_WEEKS} 週，觀察使用完成率", RED),
    ("致命", "其他設計公司有一樣的痛點，而且願意付費嗎？",
     f"{INTERVIEW_COUNT} 家訪談加上 {BETA_COUNT} 家測試", RED),
    ("高", "含/不含機制真的能減少客戶追加糾紛嗎？",
     f"追蹤 {TEST_CASES} 個案子的追加金額，和歷史平均比較", ORANGE),
    ("中", "風險報告能讓覆核時間明顯縮短嗎？",
     f"賀云實測覆核 {TEST_REPORTS} 份報告，計時比較", GOLD),
]
for i, (lv, hyp, method, c) in enumerate(risks):
    y = Inches(2.5) + Inches(i * 1.05)
    card(s, Inches(0.8), y, Inches(11.7), Inches(0.85), DARK)
    card(s, Inches(1.0), y + Inches(0.15), Inches(0.9), Inches(0.5), c)
    tx(s, Inches(1.0), y + Inches(0.18), Inches(0.9), Inches(0.45),
       lv, sz=12, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tx(s, Inches(2.2), y + Inches(0.13), Inches(5), Inches(0.55),
       hyp, sz=15, b=True)
    tx(s, Inches(7.5), y + Inches(0.13), Inches(4.7), Inches(0.55),
       method, sz=13, c=MID_GRAY)

card(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.75), DARK_GREEN_BG)
tx(s, Inches(1.2), Inches(6.4), Inches(11), Inches(0.5),
   f"初步實測：覆核時間從 {REVIEW_BEFORE} 分鐘降至 {REVIEW_AFTER} 分鐘（賀云測試 {TEST_REPORTS} 份報告）",
   sz=15, c=ACCENT, b=True)


# ════════════════════════════════════════
# 12 — ASK
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, Inches(0.8), Inches(1.8), Inches(1.2), Inches(0.08))
tx(s, Inches(0.8), Inches(2.1), Inches(11), Inches(1),
   "下一步", sz=48, b=True, c=ACCENT)

tf = tx(s, Inches(0.8), Inches(3.3), Inches(11), Inches(3), "", sz=22)
ap(tf, f"1.  最小可行產品預計 {MVP_WEEKS} 週內上線", sz=22, c=WHITE, sp=Pt(16))
ap(tf, "2.  賀云公司內部用真實案件測試", sz=22, c=WHITE, sp=Pt(16))
ap(tf, f"3.  徵求 {BETA_COUNT} 家設計公司加入測試", sz=22, c=WHITE, sp=Pt(16))
ap(tf, "4.  希望業師引薦設計公司老闆或產業資源", sz=22, c=ACCENT, b=True, sp=Pt(16))

tx(s, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
   "Q-Check — 室內裝修公司的風險管理決策系統", sz=20, c=MID_GRAY)

tx(s, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
   f"第六組：呂賀云  /  高克瑋  ·  {EMAIL}", sz=14, c=MID_GRAY)


# ── Save ──
out = "/Users/stephen/Interior_Javis/Q-Check-Pitch-Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
