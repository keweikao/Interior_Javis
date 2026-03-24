#!/usr/bin/env python3
"""
Q-Check 競爭定位圖 — 單張投影片
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Palette (warm design) ──
BG = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2B, 0x25, 0x22)
TEXT_MID = RGBColor(0x8A, 0x82, 0x7A)
AXIS_COLOR = RGBColor(0xD0, 0xC8, 0xC0)
AMBER = RGBColor(0xEF, 0x9F, 0x27)
AMBER_LIGHT = RGBColor(0xFD, 0xF3, 0xE0)
TEAL = RGBColor(0x5D, 0xCA, 0xA5)
TEAL_LIGHT = RGBColor(0xE8, 0xF8, 0xF0)
PURPLE = RGBColor(0xAF, 0xA9, 0xEC)
PURPLE_LIGHT = RGBColor(0xF0, 0xEE, 0xFB)
SKY = RGBColor(0x85, 0xB7, 0xEB)
SKY_LIGHT = RGBColor(0xEB, 0xF3, 0xFB)
GRAY_BORDER = RGBColor(0xC8, 0xC2, 0xBA)
GRAY_LIGHT = RGBColor(0xF5, 0xF3, 0xF0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.background.fill
bg.solid()
bg.fore_color.rgb = BG

W = 13.333
H = 7.5

# ── Title ──
txBox = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Q-Check 競爭格局：每個玩家賣的是什麼承諾？"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = TEXT_DARK
p.font.name = "Microsoft JhengHei"

# ── Chart area ──
# The chart area for positioning cards
CHART_L = 1.8   # inches from left
CHART_R = 11.8
CHART_T = 1.4   # top (Y=100% = event prevention)
CHART_B = 6.8   # bottom (Y=0%)
CHART_CX = (CHART_L + CHART_R) / 2
CHART_CY = (CHART_T + CHART_B) / 2

def chart_x(pct):
    """pct: 0=left, 100=right"""
    return CHART_L + (CHART_R - CHART_L) * pct / 100

def chart_y(pct):
    """pct: 0=bottom, 100=top"""
    return CHART_B - (CHART_B - CHART_T) * pct / 100


# ── Axes (thin cross lines) ──
# Vertical axis
vline = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(CHART_CX) - Pt(0.5), Inches(CHART_T), Pt(1), Inches(CHART_B - CHART_T))
vline.fill.solid()
vline.fill.fore_color.rgb = AXIS_COLOR
vline.line.fill.background()

# Horizontal axis
hline = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(CHART_L), Inches(CHART_CY) - Pt(0.5), Inches(CHART_R - CHART_L), Pt(1))
hline.fill.solid()
hline.fill.fore_color.rgb = AXIS_COLOR
hline.line.fill.background()

# ── Axis labels ──
def axis_label(x, y, text, sz=11, align=PP_ALIGN.CENTER):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(2), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = TEXT_MID
    p.font.name = "Microsoft JhengHei"
    p.alignment = align

# X axis
axis_label(CHART_L - 0.3, CHART_CY + 0.05, "通用工具", align=PP_ALIGN.LEFT)
axis_label(CHART_R - 1.7, CHART_CY + 0.05, "室內設計專用", align=PP_ALIGN.RIGHT)
# Y axis
axis_label(CHART_CX - 1.0, CHART_T - 0.35, "事前預防", align=PP_ALIGN.CENTER)
axis_label(CHART_CX - 1.0, CHART_B + 0.02, "事後管理", align=PP_ALIGN.CENTER)


# ── Card helper ──
def add_card(x_pct, y_pct, name, promise, detail,
             border_color=GRAY_BORDER, fill_color=GRAY_LIGHT,
             card_w=2.2, card_h=1.15, border_w=Pt(1), name_sz=13, promise_sz=10, detail_sz=9):
    cx = chart_x(x_pct) - card_w / 2
    cy = chart_y(y_pct) - card_h / 2

    shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(cx), Inches(cy), Inches(card_w), Inches(card_h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_w

    # Name
    tb = s.shapes.add_textbox(
        Inches(cx + 0.12), Inches(cy + 0.08), Inches(card_w - 0.24), Inches(0.3))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(name_sz)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.font.name = "Microsoft JhengHei"

    # Promise (italic with quotes)
    tb2 = s.shapes.add_textbox(
        Inches(cx + 0.12), Inches(cy + 0.38), Inches(card_w - 0.24), Inches(0.3))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = f"「{promise}」"
    p2.font.size = Pt(promise_sz)
    p2.font.italic = True
    p2.font.color.rgb = border_color if border_color != GRAY_BORDER else TEXT_MID
    p2.font.name = "Microsoft JhengHei"

    # Detail
    if detail:
        tb3 = s.shapes.add_textbox(
            Inches(cx + 0.12), Inches(cy + 0.7), Inches(card_w - 0.24), Inches(0.4))
        tf3 = tb3.text_frame; tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = detail
        p3.font.size = Pt(detail_sz)
        p3.font.color.rgb = TEXT_MID
        p3.font.name = "Microsoft JhengHei"

    return cx, cy, card_w, card_h


# ── Cards ──

# 1. Q-Check (right-upper, largest, amber)
qx, qy, qw, qh = add_card(
    78, 82, "Q-Check", "這份報價有沒有埋地雷？",
    "風險檢核 · 含/不含 · 專業判斷",
    border_color=AMBER, fill_color=AMBER_LIGHT,
    card_w=2.6, card_h=1.25, border_w=Pt(2.5), name_sz=15, promise_sz=11, detail_sz=10)

# Dashed highlight box around Q-Check
dash_margin = 0.15
dash_shape = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(qx - dash_margin), Inches(qy - dash_margin),
    Inches(qw + dash_margin * 2), Inches(qh + dash_margin * 2 + 0.35))
dash_shape.fill.background()
dash_shape.line.color.rgb = AMBER
dash_shape.line.width = Pt(1.5)
# Set dash style via XML
from lxml import etree
ln = dash_shape.line._ln
prstDash = etree.SubElement(ln, qn('a:prstDash'))
prstDash.set('val', 'dash')

# Label below dashed box
tb = s.shapes.add_textbox(
    Inches(qx - dash_margin), Inches(qy + qh + 0.05),
    Inches(qw + dash_margin * 2), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "唯一的事前預防 × 設計專用"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = AMBER
p.font.name = "Microsoft JhengHei"
p.alignment = PP_ALIGN.CENTER

# 2. 易裝修 EZid (right-lower)
add_card(73, 42, "易裝修 EZid", "案子跑更快、團隊不亂",
         "$1,200/月起 · 報價生成 · 流程管理 · 收支",
         border_color=TEAL, fill_color=TEAL_LIGHT)

# 3. Excel 報價模板 (left-center)
add_card(25, 52, "Excel 報價模板", "我最懂自己的案子",
         "免費 · 70-80% 設計師在用 · 零檢核零傳承")

# 4. LINE 群組 (left-lower)
add_card(25, 30, "LINE 群組", "傳個訊息就搞定了",
         "免費 · 90% 在用 · 即時但無法追溯")

# 5. 100室內設計 (right-bottom)
add_card(73, 18, "100室內設計", "幫你找到下一個客戶",
         "B2C 媒合平台 · 2,000+ 設計公司",
         border_color=PURPLE, fill_color=PURPLE_LIGHT)

# 6. 元欣估價系統 (left-upper)
add_card(25, 73, "元欣估價系統", "估價標準化不出錯",
         "單機版 · 通用工程 · 不懂裝修工法",
         border_color=SKY, fill_color=SKY_LIGHT)

# 7. Notion / Trello (left-bottom-corner, smallest)
add_card(20, 12, "Notion / Trello", "什麼都能管 = 什麼都不深", "",
         card_w=1.9, card_h=0.85, name_sz=11, promise_sz=9)


# ── Legend (bottom-right) ──
legend_x = 10.3
legend_y = 6.5
legend_items = [
    (AMBER, "風險預防"),
    (TEAL, "流程加速"),
    (SKY, "通用工具"),
    (GRAY_BORDER, "替代方案"),
]

for i, (color, label_text) in enumerate(legend_items):
    lx = legend_x + (i * 1.5)
    # small square
    sq = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(lx), Inches(legend_y), Inches(0.2), Inches(0.2))
    sq.fill.solid()
    sq.fill.fore_color.rgb = color
    sq.line.fill.background()
    # label
    tb = s.shapes.add_textbox(
        Inches(lx + 0.25), Inches(legend_y - 0.02), Inches(1.2), Inches(0.25))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = label_text
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MID
    p.font.name = "Microsoft JhengHei"


# ── Save ──
out = "/Users/stephen/Interior_Javis/Q-Check_競爭定位圖.pptx"
prs.save(out)
print(f"Saved: {out}")
